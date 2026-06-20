"""
Generate bilingual Vi/En presentation slides (.pptx).
Output: docs/slides/presentation.pptx
Requires: pip install python-pptx
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

SLIDES: list[tuple[str, str]] = [
    (
        "Vietnamese MCQA System\nHe thong MCQA Tieng Viet",
        "HackAIthon 2026 -- Bang C (INNOVATOR)",
    ),
    (
        "Problem / Bai toan",
        "EN: Answer Vietnamese multiple-choice questions using small "
        "open-source LLMs.\n"
        "VI: Tra loi cau hoi trac nghiem TV dung LLM nho nguon mo.\n\n"
        "Input: /data/*.csv or *.json (qid, question, choices)\n"
        "Output: /output/pred.csv (qid, answer A-K)",
    ),
    (
        "Core Method / Phuong phap chinh",
        "Constrained Likelihood Scoring\n"
        "- 1 forward pass per question\n"
        "- Extract logits for A..K at last position\n"
        "- Argmax = answer (always valid, no parsing errors)\n"
        "Tinh diem log-XS rang buoc: 1 forward pass, lay logits A..K, "
        "argmax = dap an",
    ),
    (
        "Architecture / Kien truc",
        "config.py  -> io_utils.py  -> prompt.py\n"
        "     |-> backends/stub (test)\n"
        "     |-> backends/hf   (CPU dev)\n"
        "     |-> backends/vllm (GPU)\n"
        "     |-> backends/unslo (Unsloth 4-bit)\n"
        "scorer.py -> run.py (Docker entrypoint)",
    ),
    (
        "Optimization / Toi uu hoa",
        "- AWQ 4-bit quantization (vLLM) -- luong hoa 4-bit\n"
        "- BnB 4-bit quantization (Unsloth) -- luong hoa Unsloth\n"
        "- Batch inference -- xu ly theo lo\n"
        "- Middle-truncation for long passages -- cat giua doan van dai\n"
        "- Context budget 8000 chars -- ngan sach 8000 ky tu\n"
        "- Qwen3.5-4B-Instruct: fast + accurate\n"
        "- Unsloth: memory efficient, fast kernels",
    ),
    (
        "Results / Ket qua",
        "(To be filled after benchmarking)\n\n"
        "Target: >75% accuracy on private test\n"
        "Target: >20 Q/min throughput on A100",
    ),
    (
        "Thank you / Cam on",
        "Team: ...\nContact: ...\nGitHub: ...",
    ),
]


def build_slides(output_path: str = "docs/slides/presentation.pptx"):
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # title + content

    for title_text, body_text in SLIDES:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = title_text
        tf = slide.placeholders[1].text_frame
        tf.text = body_text
        tf.word_wrap = True

    prs.save(output_path)
    print(f"Slides written to {output_path}")


if __name__ == "__main__":
    build_slides()
